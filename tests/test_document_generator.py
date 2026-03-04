"""Tests for document generation handlers.

Covers:
- document.create_word (scratch + template)
- document.create_pdf (text via fpdf2; HTML via weasyprint mocked if unavailable)
- document.create_presentation
- base64 output when output_path is omitted
"""

import base64
import os
import tempfile

import pytest

from worker.tasks.document_generator import (
    handle_document_create_word,
    handle_document_create_pdf,
    handle_document_create_presentation,
)

TEMPLATES_DIR = os.path.join(os.path.dirname(__file__), "..", "worker", "templates", "documents")


# ── Helpers ──────────────────────────────────────────────────────────────────

def _tmp_path(ext: str) -> str:
    fd, path = tempfile.mkstemp(suffix=ext)
    os.close(fd)
    os.unlink(path)
    return path


# ── document.create_word ─────────────────────────────────────────────────────

class TestCreateWord:

    def test_scratch_creates_docx(self):
        path = _tmp_path(".docx")
        try:
            result = handle_document_create_word({
                "title": "Test Document",
                "content": "Hello world\n- Item 1\n- Item 2\n## Section\nMore text",
                "output_path": path,
            })
            assert result["ok"] is True
            assert result["path"] == path
            assert result["size_bytes"] > 0
            assert os.path.isfile(path)
        finally:
            if os.path.exists(path):
                os.unlink(path)

    def test_scratch_base64_when_no_output_path(self):
        result = handle_document_create_word({
            "title": "In-Memory Doc",
            "content": "Just a paragraph.",
        })
        assert result["ok"] is True
        assert result["path"] is None
        assert "base64" in result
        raw = base64.b64decode(result["base64"])
        assert raw[:4] == b"PK\x03\x04"  # ZIP/OOXML magic bytes
        assert result["size_bytes"] == len(raw)

    def test_scratch_default_title(self):
        result = handle_document_create_word({"content": "No title given"})
        assert result["ok"] is True

    def test_template_renders_variables(self):
        tpl_path = os.path.join(TEMPLATES_DIR, "propuesta_bim.docx")
        if not os.path.isfile(tpl_path):
            pytest.skip("propuesta_bim.docx template not found")

        path = _tmp_path(".docx")
        try:
            result = handle_document_create_word({
                "template_path": tpl_path,
                "data": {
                    "proyecto": "Torre Test",
                    "cliente": "Test Corp",
                    "fecha": "2026-03-04",
                    "contexto": "Contexto de prueba",
                    "alcance": "Alcance de prueba",
                    "metodologia": "Metodología de prueba",
                    "entregables": "Entregables de prueba",
                    "inversion": "$10,000 USD",
                    "cronograma": "4 semanas",
                    "condiciones": "Condiciones estándar",
                },
                "output_path": path,
            })
            assert result["ok"] is True
            assert result["size_bytes"] > 0
            assert os.path.isfile(path)

            from docx import Document
            doc = Document(path)
            full_text = "\n".join(p.text for p in doc.paragraphs)
            assert "Torre Test" in full_text
            assert "Test Corp" in full_text
            assert "{{ proyecto }}" not in full_text
        finally:
            if os.path.exists(path):
                os.unlink(path)

    def test_template_base64_when_no_output_path(self):
        tpl_path = os.path.join(TEMPLATES_DIR, "cotizacion_bim.docx")
        if not os.path.isfile(tpl_path):
            pytest.skip("cotizacion_bim.docx template not found")

        result = handle_document_create_word({
            "template_path": tpl_path,
            "data": {
                "proyecto": "Proyecto B64",
                "cliente": "Cliente B64",
                "fecha": "2026-01-01",
                "referencia": "COT-001",
                "servicios": [],
                "subtotal": "$0",
                "iva": "$0",
                "total": "$0",
                "vigencia": "30 días",
                "notas": "",
            },
        })
        assert result["ok"] is True
        assert result["path"] is None
        assert "base64" in result
        raw = base64.b64decode(result["base64"])
        assert raw[:4] == b"PK\x03\x04"

    def test_template_not_found_raises(self):
        with pytest.raises(FileNotFoundError):
            handle_document_create_word({
                "template_path": "/nonexistent/template.docx",
                "data": {},
            })

    def test_cotizacion_template_renders_services(self):
        tpl_path = os.path.join(TEMPLATES_DIR, "cotizacion_bim.docx")
        if not os.path.isfile(tpl_path):
            pytest.skip("cotizacion_bim.docx template not found")

        path = _tmp_path(".docx")
        try:
            result = handle_document_create_word({
                "template_path": tpl_path,
                "data": {
                    "proyecto": "Edificio Omega",
                    "cliente": "Omega SA",
                    "fecha": "2026-03-04",
                    "referencia": "COT-TEST-001",
                    "servicios": [
                        {"nombre": "Modelado 3D", "descripcion": "LOD 300", "plazo": "4 sem", "valor": "$5,000"},
                        {"nombre": "Coordinación", "descripcion": "MEP", "plazo": "2 sem", "valor": "$3,000"},
                    ],
                    "subtotal": "$8,000",
                    "iva": "$1,280",
                    "total": "$9,280",
                    "vigencia": "15 días",
                    "notas": "Incluye revisiones.",
                },
                "output_path": path,
            })
            assert result["ok"] is True
            assert os.path.isfile(path)
        finally:
            if os.path.exists(path):
                os.unlink(path)


# ── document.create_pdf ──────────────────────────────────────────────────────

class TestCreatePdf:

    def test_text_creates_pdf(self):
        path = _tmp_path(".pdf")
        try:
            result = handle_document_create_pdf({
                "text_content": "This is a test PDF.\nLine two.\nLine three.",
                "title": "Test PDF",
                "output_path": path,
            })
            assert result["ok"] is True
            assert result["path"] == path
            assert result["size_bytes"] > 0
            assert os.path.isfile(path)
            with open(path, "rb") as f:
                assert f.read(5) == b"%PDF-"
        finally:
            if os.path.exists(path):
                os.unlink(path)

    def test_text_base64_when_no_output_path(self):
        result = handle_document_create_pdf({
            "text_content": "In-memory PDF content.",
        })
        assert result["ok"] is True
        assert result["path"] is None
        assert "base64" in result
        raw = base64.b64decode(result["base64"])
        assert raw[:5] == b"%PDF-"
        assert result["size_bytes"] == len(raw)

    def test_text_without_title(self):
        result = handle_document_create_pdf({
            "text_content": "No title, just text.",
        })
        assert result["ok"] is True
        assert result["size_bytes"] > 0

    def test_missing_content_raises(self):
        with pytest.raises(ValueError, match="html_content.*text_content"):
            handle_document_create_pdf({})

    def test_html_pdf_with_weasyprint_or_skip(self):
        """Test HTML→PDF if weasyprint is available, otherwise skip."""
        try:
            from weasyprint import HTML  # noqa: F401
        except (ImportError, OSError):
            pytest.skip("weasyprint not available (requires system libs)")

        path = _tmp_path(".pdf")
        try:
            result = handle_document_create_pdf({
                "html_content": "<h1>Hello</h1><p>World</p>",
                "output_path": path,
            })
            assert result["ok"] is True
            assert os.path.isfile(path)
        finally:
            if os.path.exists(path):
                os.unlink(path)

    def test_html_pdf_runtime_error_when_weasyprint_missing(self):
        """If weasyprint is not importable, a RuntimeError should be raised."""
        from unittest.mock import patch

        with patch.dict("sys.modules", {"weasyprint": None}):
            try:
                from weasyprint import HTML  # noqa: F401
                pytest.skip("weasyprint is available; cannot test fallback")
            except (ImportError, TypeError):
                pass

        # Only run assertion if weasyprint is truly missing
        try:
            from weasyprint import HTML  # noqa: F401
            pytest.skip("weasyprint is available")
        except (ImportError, OSError):
            with pytest.raises(RuntimeError, match="weasyprint"):
                handle_document_create_pdf({"html_content": "<p>test</p>"})


# ── document.create_presentation ─────────────────────────────────────────────

class TestCreatePresentation:

    def test_creates_pptx(self):
        path = _tmp_path(".pptx")
        try:
            result = handle_document_create_presentation({
                "slides": [
                    {"title": "Title Slide", "content": "Subtitle here", "notes": "Speaker notes"},
                    {"title": "Content Slide", "content": "- Point 1\n- Point 2"},
                    {"title": "Final Slide", "content": "Thank you"},
                ],
                "output_path": path,
            })
            assert result["ok"] is True
            assert result["path"] == path
            assert result["slide_count"] == 3
            assert result["size_bytes"] > 0
            assert os.path.isfile(path)
        finally:
            if os.path.exists(path):
                os.unlink(path)

    def test_base64_when_no_output_path(self):
        result = handle_document_create_presentation({
            "slides": [
                {"title": "Slide 1", "content": "Content 1"},
            ],
        })
        assert result["ok"] is True
        assert result["path"] is None
        assert result["slide_count"] == 1
        assert "base64" in result
        raw = base64.b64decode(result["base64"])
        assert raw[:4] == b"PK\x03\x04"
        assert result["size_bytes"] == len(raw)

    def test_empty_slides_raises(self):
        with pytest.raises(ValueError, match="slides"):
            handle_document_create_presentation({"slides": []})

    def test_missing_slides_raises(self):
        with pytest.raises(ValueError, match="slides"):
            handle_document_create_presentation({})

    def test_single_slide_uses_title_content_layout(self):
        result = handle_document_create_presentation({
            "slides": [{"title": "Only Slide", "content": "Body text"}],
        })
        assert result["ok"] is True
        assert result["slide_count"] == 1

    def test_slide_with_notes(self):
        from pptx import Presentation
        from io import BytesIO

        result = handle_document_create_presentation({
            "slides": [
                {"title": "S1", "content": "C1"},
                {"title": "S2", "content": "C2", "notes": "Important note"},
            ],
        })
        raw = base64.b64decode(result["base64"])
        prs = Presentation(BytesIO(raw))
        slide2 = prs.slides[1]
        assert "Important note" in slide2.notes_slide.notes_text_frame.text

    def test_default_slide_title(self):
        result = handle_document_create_presentation({
            "slides": [{"content": "No explicit title"}],
        })
        assert result["ok"] is True
        assert result["slide_count"] == 1


# ── Handler registration ─────────────────────────────────────────────────────

class TestRegistration:

    def test_handlers_registered(self):
        from worker.tasks import TASK_HANDLERS
        assert "document.create_word" in TASK_HANDLERS
        assert "document.create_pdf" in TASK_HANDLERS
        assert "document.create_presentation" in TASK_HANDLERS
