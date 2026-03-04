"""Task handlers: document.create_word, document.create_pdf, document.create_presentation.

Generate professional documents (Word, PDF, PowerPoint) from templates or from scratch.
"""

import base64
import logging
import os
from io import BytesIO
from typing import Any, Dict

logger = logging.getLogger("worker.tasks.document_generator")


# ---------------------------------------------------------------------------
# document.create_word
# ---------------------------------------------------------------------------

def handle_document_create_word(input_data: Dict[str, Any]) -> Dict[str, Any]:
    """Create a Word (.docx) document.

    With template_path + data  → renders a .docx template via docxtpl (Jinja2).
    Without template           → builds a document from scratch via python-docx.

    Input (template mode):
        template_path: str   — path to a .docx template with Jinja2 tags
        data: dict           — variables to render
        output_path: str     — (optional) where to save; omit for base64 response

    Input (scratch mode):
        title: str           — document title
        content: str         — body text (plain text; lines starting with "- " become bullets)
        output_path: str     — (optional)

    Returns:
        ok: bool, path: str | None, size_bytes: int, [base64: str]
    """
    template_path = input_data.get("template_path")
    output_path = input_data.get("output_path")

    if template_path:
        return _word_from_template(template_path, input_data.get("data", {}), output_path)
    return _word_from_scratch(input_data, output_path)


def _word_from_template(template_path: str, data: dict, output_path: str | None) -> dict:
    from docxtpl import DocxTemplate

    if not os.path.isfile(template_path):
        raise FileNotFoundError(f"Template not found: {template_path}")

    tpl = DocxTemplate(template_path)
    tpl.render(data)

    if output_path:
        os.makedirs(os.path.dirname(output_path) or ".", exist_ok=True)
        tpl.save(output_path)
        size = os.path.getsize(output_path)
        logger.info("Word (template) saved: %s (%d bytes)", output_path, size)
        return {"ok": True, "path": output_path, "size_bytes": size}

    buf = BytesIO()
    tpl.save(buf)
    b64 = base64.b64encode(buf.getvalue()).decode()
    logger.info("Word (template) generated in-memory (%d bytes)", len(buf.getvalue()))
    return {"ok": True, "path": None, "size_bytes": len(buf.getvalue()), "base64": b64}


def _word_from_scratch(input_data: dict, output_path: str | None) -> dict:
    from docx import Document

    title = input_data.get("title", "Documento")
    content = input_data.get("content", "")

    doc = Document()
    doc.add_heading(title, level=0)

    for line in content.split("\n"):
        stripped = line.strip()
        if not stripped:
            continue
        if stripped.startswith("## "):
            doc.add_heading(stripped[3:], level=2)
        elif stripped.startswith("# "):
            doc.add_heading(stripped[2:], level=1)
        elif stripped.startswith("- ") or stripped.startswith("* "):
            doc.add_paragraph(stripped[2:], style="List Bullet")
        else:
            doc.add_paragraph(stripped)

    if output_path:
        os.makedirs(os.path.dirname(output_path) or ".", exist_ok=True)
        doc.save(output_path)
        size = os.path.getsize(output_path)
        logger.info("Word (scratch) saved: %s (%d bytes)", output_path, size)
        return {"ok": True, "path": output_path, "size_bytes": size}

    buf = BytesIO()
    doc.save(buf)
    b64 = base64.b64encode(buf.getvalue()).decode()
    logger.info("Word (scratch) generated in-memory (%d bytes)", len(buf.getvalue()))
    return {"ok": True, "path": None, "size_bytes": len(buf.getvalue()), "base64": b64}


# ---------------------------------------------------------------------------
# document.create_pdf
# ---------------------------------------------------------------------------

def handle_document_create_pdf(input_data: Dict[str, Any]) -> Dict[str, Any]:
    """Create a PDF document.

    With html_content  → renders via weasyprint (requires system libs: Cairo, Pango).
    With text_content  → renders via fpdf2 (pure Python, no system deps).

    Input (HTML mode):
        html_content: str  — HTML string
        output_path: str   — (optional)

    Input (text mode):
        text_content: str  — plain text
        title: str         — (optional) PDF title
        output_path: str   — (optional)

    Returns:
        ok: bool, path: str | None, size_bytes: int, [base64: str]
    """
    html_content = input_data.get("html_content")
    text_content = input_data.get("text_content")
    output_path = input_data.get("output_path")

    if html_content:
        return _pdf_from_html(html_content, output_path)
    if text_content:
        return _pdf_from_text(text_content, input_data.get("title"), output_path)

    raise ValueError("Either 'html_content' or 'text_content' is required")


def _pdf_from_html(html_content: str, output_path: str | None) -> dict:
    try:
        from weasyprint import HTML  # type: ignore[import-untyped]
    except (ImportError, OSError) as exc:
        raise RuntimeError(
            "weasyprint is not available (requires system libs: Cairo, Pango). "
            f"Falling back is not supported for HTML PDFs. Error: {exc}"
        ) from exc

    if output_path:
        os.makedirs(os.path.dirname(output_path) or ".", exist_ok=True)
        HTML(string=html_content).write_pdf(output_path)
        size = os.path.getsize(output_path)
        logger.info("PDF (HTML) saved: %s (%d bytes)", output_path, size)
        return {"ok": True, "path": output_path, "size_bytes": size}

    pdf_bytes = HTML(string=html_content).write_pdf()
    b64 = base64.b64encode(pdf_bytes).decode()
    logger.info("PDF (HTML) generated in-memory (%d bytes)", len(pdf_bytes))
    return {"ok": True, "path": None, "size_bytes": len(pdf_bytes), "base64": b64}


def _pdf_from_text(text_content: str, title: str | None, output_path: str | None) -> dict:
    from fpdf import FPDF

    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()

    if title:
        pdf.set_font("helvetica", "B", size=18)
        pdf.cell(w=0, h=12, text=title)
        pdf.ln(14)

    pdf.set_font("helvetica", size=12)
    for line in text_content.split("\n"):
        if not line.strip():
            pdf.ln(4)
        else:
            pdf.multi_cell(w=0, h=6, text=line)
            pdf.ln(1)

    if output_path:
        os.makedirs(os.path.dirname(output_path) or ".", exist_ok=True)
        pdf.output(output_path)
        size = os.path.getsize(output_path)
        logger.info("PDF (text) saved: %s (%d bytes)", output_path, size)
        return {"ok": True, "path": output_path, "size_bytes": size}

    pdf_bytes = pdf.output()
    b64 = base64.b64encode(pdf_bytes).decode()
    logger.info("PDF (text) generated in-memory (%d bytes)", len(pdf_bytes))
    return {"ok": True, "path": None, "size_bytes": len(pdf_bytes), "base64": b64}


# ---------------------------------------------------------------------------
# document.create_presentation
# ---------------------------------------------------------------------------

def handle_document_create_presentation(input_data: Dict[str, Any]) -> Dict[str, Any]:
    """Create a PowerPoint (.pptx) presentation.

    Input:
        slides: list[dict]  — each with 'title', 'content', optional 'notes'
        output_path: str    — (optional)

    Returns:
        ok: bool, path: str | None, slide_count: int, size_bytes: int, [base64: str]
    """
    from pptx import Presentation

    slides_data = input_data.get("slides")
    if not slides_data or not isinstance(slides_data, list):
        raise ValueError("'slides' is required and must be a non-empty list")

    output_path = input_data.get("output_path")

    prs = Presentation()

    for i, sd in enumerate(slides_data):
        title = sd.get("title", f"Slide {i + 1}")
        content = sd.get("content", "")
        notes = sd.get("notes")

        if i == 0 and len(slides_data) > 1:
            layout = prs.slide_layouts[0]  # Title Slide
        else:
            layout = prs.slide_layouts[1]  # Title + Content

        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = title

        if layout == prs.slide_layouts[0]:
            if 1 in {p.placeholder_format.idx for p in slide.placeholders}:
                slide.placeholders[1].text = content
        else:
            slide.placeholders[1].text = content

        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    if output_path:
        os.makedirs(os.path.dirname(output_path) or ".", exist_ok=True)
        prs.save(output_path)
        size = os.path.getsize(output_path)
        logger.info("PPTX saved: %s (%d bytes, %d slides)", output_path, size, len(slides_data))
        return {"ok": True, "path": output_path, "slide_count": len(slides_data), "size_bytes": size}

    buf = BytesIO()
    prs.save(buf)
    b64 = base64.b64encode(buf.getvalue()).decode()
    logger.info("PPTX generated in-memory (%d bytes, %d slides)", len(buf.getvalue()), len(slides_data))
    return {
        "ok": True,
        "path": None,
        "slide_count": len(slides_data),
        "size_bytes": len(buf.getvalue()),
        "base64": b64,
    }
